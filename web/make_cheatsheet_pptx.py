# Build an editable PowerPoint version of the Liu cheat-sheet poster
# (docs/liu-cheatsheet.pptx). One A4 portrait slide; every text element is
# a real text box, cards and code panels are shapes, figures are embedded
# PNGs (real interpreter output).
#
# Usage: python3 web/make_cheatsheet_pptx.py <dir-with-fig_{fm,diff,svgd}.png>
# The PNGs are the three <svg> figures of web/cheatsheet-poster.html,
# rasterized at 3x (any SVG renderer works).
import os, sys
FIGDIR = sys.argv[1] if len(sys.argv) > 1 else '.'
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

INK    = RGBColor(0x10, 0x18, 0x20)   # page background
CARD   = RGBColor(0x16, 0x21, 0x2B)   # card background
PANEL  = RGBColor(0x1D, 0x2A, 0x36)   # code panel
BORDER = RGBColor(0x2A, 0x39, 0x48)
TEXT   = RGBColor(0xE7, 0xED, 0xF2)
MUTED  = RGBColor(0x8F, 0xA0, 0xAE)
CODE   = RGBColor(0xDB, 0xE5, 0xEC)
CYAN   = RGBColor(0x4F, 0xC3, 0xD9)
SAND   = RGBColor(0xD9, 0xA9, 0x6C)

MONO, SANS = "Consolas", "Segoe UI"

prs = Presentation()
prs.slide_width, prs.slide_height = Cm(21.0), Cm(29.7)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

def bg():
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Cm(21.0), Cm(29.7))
    r.fill.solid(); r.fill.fore_color.rgb = INK
    r.line.fill.background()
    r.shadow.inherit = False
    return r

def box(x, y, w, h, fill=None, line=None, radius=None):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    r = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if radius is not None:
        try:
            r.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        r.fill.background()
    else:
        r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(0.75)
    r.shadow.inherit = False
    return r

def text(x, y, w, h, runs_lines, size=10, font=SANS, color=TEXT, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None, wrap=True):
    """runs_lines: list of paragraphs; each is a list of (text, overrides) runs."""
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        for t, ov in line:
            r = p.add_run(); r.text = t
            f = r.font
            f.size = Pt(ov.get('size', size))
            f.name = ov.get('font', font)
            f.color.rgb = ov.get('color', color)
            f.bold = ov.get('bold', bold)
    return tb

def set_ea_font(tb, name):   # East-Asian font for the 流 glyph
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {}); rPr.append(ea)
            ea.set('typeface', name)

bg()

ML, MR, MT = 0.9, 0.9, 0.8          # page margins (cm)
W = 21.0 - ML - MR                  # 19.2

# ---- header -------------------------------------------------------------
mark = text(ML, MT, 4.2, 1.6,
    [[("流", {'size': 30, 'bold': True, 'color': TEXT}),
      ("  Liu", {'size': 13, 'bold': True, 'color': CYAN})]])
set_ea_font(mark, "Microsoft YaHei")

text(ML + 3.4, MT + 0.05, 8.5, 2.6, [
    [("An executable notation for measure transport.", {'bold': True, 'color': TEXT})],
    [("Straight-line programs: no loops, no I/O, not Turing complete.", {})],
    [("Text + ", {}), ("seed", {'font': MONO, 'color': CYAN}),
     (" determine every output ", {}), ("bit for bit", {'bold': True, 'color': TEXT}), (".", {})],
], size=10, color=MUTED, line_spacing=1.25)

pw = 7.0
box(21.0 - MR - pw, MT, pw, 1.75, fill=CARD, line=BORDER, radius=0.25)
text(21.0 - MR - pw, MT + 0.22, pw, 1.4, [
    [("flow( field( path ) ) # μ", {'size': 12.5, 'font': MONO, 'color': CYAN})],
    [("THE ONE PIPELINE", {'size': 8, 'color': MUTED})],
], align=PP_ALIGN.CENTER)

# ---- cards --------------------------------------------------------------
CARD_TOP, CARD_H, GAP = 4.05, 7.95, 0.35
FIG_W = 6.6
PAD = 0.45

cards = [
  dict(title="Flow Matching", em="· declared path",
    desc=[[("The interpolation formula defines a path. ", {}),
           ("field", {'font': MONO, 'color': CYAN}),
           (" regresses E[ẋ", {}), ("t", {'size': 7}), (" | x", {}), ("t", {'size': 7}),
           ("] — the L² minimizer, solving the continuity equation for this path. ", {}),
           ("flow", {'font': MONO, 'color': CYAN}),
           (" defines a pushforward using the field.", {})]],
    code=["xt = t*data + (1-t)*noise",
          "pt = prob(xt)",
          "v  = field(pt, estimator=regress(mlp(2->64->64->2)))",
          "T  = flow(v, steps=50)",
          "plot data, (T # noise) ~ 1000"],
    fig=os.path.join(FIGDIR, "fig_fm.png"), chips=[(SAND, "target sample"), (CYAN, "T # noise")]),
  dict(title="Diffusion", em="· one line away",
    desc=[[("Change the interpolation — same regression.", {})]],
    code=["xt = sqrt(1-t*t)*noise + t*data",
          "pt = prob(xt)",
          "v  = field(pt, estimator=regress(mlp(2->96->96->2)))",
          "T  = flow(v, steps=60)",
          "plot data, (T # noise) ~ 800"],
    fig=os.path.join(FIGDIR, "fig_diff.png"), chips=[(SAND, "target sample"), (CYAN, "T # noise")]),
  dict(title="SVGD", em="· gradient flow",
    desc=[[("descent", {'font': MONO, 'color': CYAN}),
           (" defines the steepest descent of a divergence in P(ℝ²). "
            "No regression — the field is lazily estimated each step.", {})]],
    code=["cloud  = uniform([-3,-3], [3,3]) ~ 300",
          "target = 0.7*gaussian([-2,0], 0.4)",
          "       or 0.3*gaussian([ 2,0], 0.4)",
          "qt = descent(reverseKL(target), from=cloud)",
          "v  = field(qt, estimator=nw(kernel=rbf))",
          "T  = flow(v, steps=400, lr=0.8)",
          "plot target ~ 1500, (T # cloud)"],
    fig=os.path.join(FIGDIR, "fig_svgd.png"), chips=[(SAND, "target sample"), (CYAN, "T # cloud")]),
]

for i, c in enumerate(cards):
    y0 = CARD_TOP + i * (CARD_H + GAP)
    box(ML, y0, W, CARD_H, fill=CARD, line=BORDER, radius=0.06)
    lx, lw = ML + PAD, W - 2*PAD - FIG_W - 0.45
    # title
    text(lx, y0 + 0.45, lw, 0.9,
        [[(c['title'], {'size': 15, 'bold': True, 'color': TEXT}),
          ("  " + c['em'], {'size': 11, 'bold': True, 'color': CYAN})]])
    # description
    text(lx, y0 + 1.35, lw, 1.35, c['desc'], size=10.5, color=MUTED, line_spacing=1.2)
    # code panel
    n = len(c['code'])
    ch = 0.56 * n + 0.42
    cy = y0 + CARD_H - PAD - ch
    box(lx, cy, lw, ch, fill=PANEL, radius=0.10)
    text(lx + 0.35, cy + 0.25, lw - 0.7, ch - 0.5,
         [[(ln, {})] for ln in c['code']],
         size=9.2, font=MONO, color=CODE, line_spacing=1.25)
    # figure
    fs = CARD_H - 2*PAD - 0.55
    fx = ML + W - PAD - FIG_W + (FIG_W - fs)/2
    slide.shapes.add_picture(c['fig'], Cm(fx), Cm(y0 + PAD), Cm(fs), Cm(fs))
    # legend
    runs = []
    for col, lbl in c['chips']:
        runs.append(("● ", {'color': col}))
        runs.append((lbl + "   ", {}))
    text(ML + W - PAD - FIG_W, y0 + PAD + fs + 0.12, FIG_W, 0.5, [runs],
         size=9, color=MUTED, align=PP_ALIGN.CENTER)

# ---- footer -------------------------------------------------------------
fy = 29.7 - 0.95
box(ML, fy - 0.15, W, 0.02, fill=BORDER)
text(ML, fy, 9.5, 0.6,
     [[("./interpreter/build.sh", {'font': MONO, 'color': CODE}),
       ("  ·  ", {}), ("python3 web/server.py", {'font': MONO, 'color': CODE})]],
     size=7, color=MUTED)
text(21.0 - MR - 11.5, fy, 11.5, 0.6,
     [[("Liu (流) v0.3+ · Juzhen (C++/BLAS) · ", {}),
       ("docs/liu-reference.md", {'font': MONO, 'color': CODE}),
       (" · figures: real interpreter output", {})]],
     size=7, color=MUTED, align=PP_ALIGN.RIGHT)

prs.save('/home/user/Juzhen/docs/liu-cheatsheet.pptx')
print('saved docs/liu-cheatsheet.pptx')
